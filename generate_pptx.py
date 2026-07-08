from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # --- HELPER FUNCTIONS ---
    def add_slide_title_content(title_text, subtitle_text=None):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        # Style Title
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204) # Brand Blue

        if subtitle_text:
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.text = subtitle_text
        return slide

    # --- SLIDE 1: TITLE ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "DECKMOUNT ELECTRONICS"
    subtitle.text = "Vision 2026: Strategic Diagnostic Hub\nAnkur Sir's Executive Roadmap"
    
    # Style
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)

    # --- SLIDE 2: THE 4 STRATEGIC QUESTIONS ---
    slide = add_slide_title_content("The Strategic Framework")
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Every product in our 2026 lineup answers four critical questions:"
    
    p = tf.add_paragraph()
    p.text = "1. What problem are we facing before?"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Why we are doing this?"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. How is it different from the market?"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "4. What is the Coming 6-Month Program?"
    p.level = 1

    # --- SLIDE 3: RHYTHM ULTRAMAX (AI-PORTABLE) ---
    slide = add_slide_title_content("Rhythm UltraMax (AI-Portable)")
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Problem: Portable ECGs compromise on signal quality (8-12 bit) and lack AI interpretation."
    tf.add_paragraph().text = "Why: To deliver hospital-grade 12-lead accuracy in a 298g handheld form factor (24-bit ADC)."
    tf.add_paragraph().text = "Difference: 24-bit ADC + 99.9% Accuracy + CardioX AI Software + BLE 5.0 Cloud Sync."
    tf.add_paragraph().text = "6-Month Program: Official Launch of CardioX Mobile App & CDSCO Scaling."

    # --- SLIDE 4: 12-CHANNEL ECG (TABLE TOP) ---
    slide = add_slide_title_content("12-Channel ECG (Table Top)")
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Problem: Slow acquisition (30s+) and manual errors delay treatment in the 'Golden Hour'."
    tf.add_paragraph().text = "Why: 10-second acquisition speed and 99.7% trace accuracy for all clinics."
    tf.add_paragraph().text = "Difference: AI-automated measurements at 1/4th the price of imports (₹90,000)."
    tf.add_paragraph().text = "6-Month Program: PHC deployment across North India + CardioX Cloud Sync."

    # --- SLIDE 5: BIPAP VT SERIES (VT400 ST) ---
    slide = add_slide_title_content("BiPAP VT Series (VT400 ST)")
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Problem: Fixed pressure discomfort and poor filtration in traditional BiPAPs."
    tf.add_paragraph().text = "Why: 'Pure Therapy' with 99.97% HEPA filtration and adaptive VGPS modes."
    tf.add_paragraph().text = "Difference: Integrated PM2.5 capture + Volume Assurance (2000ml) + Made in India."
    tf.add_paragraph().text = "6-Month Program: Launch of centralized remote monitoring dashboard."

    # --- SLIDE 6: SLEEP SENSE LEVEL 3 (SS100) ---
    slide = add_slide_title_content("Sleep Sense Level 3 (SS100)")
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Problem: Expensive hospital stays and bulky equipment for sleep studies."
    tf.add_paragraph().text = "Why: Democratize sleep diagnostics with a 214g wearable for home clinical studies."
    tf.add_paragraph().text = "Difference: 8GB internal storage + 95% automated analysis accuracy."
    tf.add_paragraph().text = "6-Month Program: Mobile app integration for real-time patient alerts."

    # --- SLIDE 7: CARDIOX SOFTWARE SUITE ---
    slide = add_slide_title_content("CardioX Software Suite: Clinical Intelligence")
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Beyond Visualization: Automated Intelligence Engine"
    tf.add_paragraph().text = "- Automated PQRS Metrics (Bazett's Correction)"
    tf.add_paragraph().text = "- Signal Processing (50/60Hz Artifact Suppression)"
    tf.add_paragraph().text = "- Interactive Tools (Virtual Digital Calipers)"
    tf.add_paragraph().text = "- AI Diagnostic Assist (Cardmia Clinical Chatbot)"

    # --- SLIDE 8: DETECTION CAPABILITIES ---
    slide = add_slide_title_content("Detection Power: 20+ ECG Findings")
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Our ECG engine detects rhythm and morphology findings with 99.9% sensitivity:"
    tf.add_paragraph().text = "- 11 Core Rhythms: VF/VT, AFib, Flutter, AV Blocks (I, II, III)"
    tf.add_paragraph().text = "- Morphology: RBBB, LBBB, PVC, PAC, ST Elevation/Depression"
    tf.add_paragraph().text = "- Reliability: Parallel Map + Serial Reduce analysis pattern"

    # --- SLIDE 9: TECHNICAL ARCHITECTURE ---
    slide = add_slide_title_content("Technical Blueprint: Holter Pipeline")
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Safe Parallelism & Fault Tolerance"
    tf.add_paragraph().text = "Dashboard -> ECG Page -> Holter Writer -> Analysis Worker -> Report -> History"
    tf.add_paragraph().text = "- Clinical Config Layer: YAML-driven thresholds (VT rates, pauses)"
    tf.add_paragraph().text = "- Session Integrity: Fault-tolerant storage domains"
    tf.add_paragraph().text = "- Replay Engine: Reconstruct sessions from metrics.jsonl"

    # --- SLIDE 10: EXECUTIVE SUMMARY & ROADMAP ---
    slide = add_slide_title_content("2026 Strategic Roadmap")
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Ankur Sir's 6-Month Goals:"
    tf.add_paragraph().text = "1. CDSCO Scaling: UltraMax & VT Series market expansion."
    tf.add_paragraph().text = "2. Mobile App Launch: Integrated CardioX App for all devices."
    tf.add_paragraph().text = "3. PHC Deployment: 500+ Clinics across North India."
    tf.add_paragraph().text = "4. Remote Monitoring: Cloud Dashboard v1.0 for hospital oversight."

    # Save the presentation
    pptx_path = "/Users/deckmount/presentationcontent/Deckmount_Strategic_Vision_2026.pptx"
    prs.save(pptx_path)
    return pptx_path

if __name__ == "__main__":
    path = create_presentation()
    print(f"PPTX generated at: {path}")
