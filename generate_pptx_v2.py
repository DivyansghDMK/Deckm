from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_professional_presentation():
    prs = Presentation()
    
    # --- COLORS ---
    BLUE = RGBColor(0, 102, 204)
    SLATE = RGBColor(15, 23, 42)
    EMERALD = RGBColor(16, 185, 129)
    WHITE = RGBColor(255, 255, 255)

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_image_to_slide(slide, img_path, left, top, width=None, height=None):
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, left, top, width=width, height=height)
        else:
            print(f"Image not found: {img_path}")

    # --- SLIDE 1: TITLE ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_background(slide, SLATE)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "DECKMOUNT ELECTRONICS"
    title.text_frame.paragraphs[0].font.color.rgb = BLUE
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    subtitle.text = "Vision 2026: Strategic Diagnostic Hub\nAnkur Sir's Executive Roadmap"
    subtitle.text_frame.paragraphs[0].font.color.rgb = WHITE
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)

    # --- SLIDE 2: THE 4 QUESTIONS ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Strategic Framework"
    tf = slide.placeholders[1].text_frame
    tf.text = "Every product in our 2026 lineup answers four critical questions:"
    for q in ["Problem Faced Before", "Why we are doing this?", "Market Differentiation", "Coming 6-Month Program"]:
        p = tf.add_paragraph()
        p.text = f"• {q}"
        p.level = 1

    # --- SLIDE 3: RHYTHM ULTRAMAX ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Rhythm UltraMax (AI-Portable)"
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "Problem: Compromised signal quality in portables."
    p = tf.add_paragraph()
    p.text = "Why: 12-lead accuracy in 298g handheld."
    p = tf.add_paragraph()
    p.text = "Edge: 24-bit ADC + AI + Cloud Sync."
    add_image_to_slide(slide, "/Users/deckmount/presentationcontent/app/src/assets/extracted/rhythm_ultramax_v3.png", Inches(5.5), Inches(1.5), width=Inches(3.5))

    # --- SLIDE 4: 12-CHANNEL ECG ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "12-Channel ECG (Table Top)"
    txBox = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.text = "Clinical Hub for PHCs"
    tf.add_paragraph().text = "• 10s Simultaneous Acquisition"
    tf.add_paragraph().text = "• 99.7% Clinical Grade Accuracy"
    tf.add_paragraph().text = "• CardioX Cloud PDF Sync"
    add_image_to_slide(slide, "/Users/deckmount/presentationcontent/app/src/assets/extracted/ecg_tabletop.jpg", Inches(0.5), Inches(1.5), width=Inches(4))

    # --- SLIDE 5: BIPAP VT SERIES ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "BiPAP VT Series (VT400 ST)"
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5), Inches(5))
    tf = txBox.text_frame
    tf.text = "Pure Therapy & Advanced Ventilation"
    tf.add_paragraph().text = "• 99.97% HEPA Filtration (PM2.5)"
    tf.add_paragraph().text = "• Adaptive VGPS / ST Modes"
    tf.add_paragraph().text = "• Volume Assurance up to 2000ml"
    add_image_to_slide(slide, "/Users/deckmount/presentationcontent/app/src/assets/extracted/bipap_new.jpg", Inches(5.5), Inches(1.5), width=Inches(3.5))

    # --- SLIDE 6: SLEEP SENSE ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Sleep Sense Level 3 (SS100)"
    txBox = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4.5), Inches(5))
    tf = txBox.text_frame
    tf.text = "Home Clinical Sleep Studies"
    tf.add_paragraph().text = "• 214g Ultra-Light Wearable"
    tf.add_paragraph().text = "• 8GB Internal Storage (16k hrs)"
    tf.add_paragraph().text = "• 95%+ AI Auto-Scoring Accuracy"
    add_image_to_slide(slide, "/Users/deckmount/presentationcontent/app/src/assets/extracted/sleepsense_new.png", Inches(0.5), Inches(1.5), width=Inches(4))

    # --- SLIDE 7: DETECTION CAPABILITIES ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, EMERALD)
    title = slide.shapes.title
    title.text = "CardioX: Sharpest Sensitivity Ever"
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Detects 30+ Critical Abnormalities"
    p = tf.add_paragraph()
    p.text = "• 12 Arrhythmias: AFib, VF/VT, Brady/Tachycardia, PVC/PAC"
    p.font.color.rgb = WHITE
    p = tf.add_paragraph()
    p.text = "• 10 Conduction: AV Blocks (I, II, III), RBBB, LBBB, WPW"
    p.font.color.rgb = WHITE
    p = tf.add_paragraph()
    p.text = "• 08 Morphology: ST Elevation/Depression, T-Wave Inversion"
    p.font.color.rgb = WHITE

    # --- SLIDE 8: ROADMAP ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, SLATE)
    title = slide.shapes.title
    title.text = "Coming 6-Month Strategic Program"
    title.text_frame.paragraphs[0].font.color.rgb = BLUE
    tf = slide.placeholders[1].text_frame
    for goal in ["CDSCO Scaling: UltraMax & VT Series", "Official CardioX Mobile App Launch", "PHC Deployment: 500+ Clinics", "Cloud Monitoring Dashboard v1.0"]:
        p = tf.add_paragraph()
        p.text = f"• {goal}"
        p.font.color.rgb = WHITE
    
    pptx_path = "/Users/deckmount/presentationcontent/Deckmount_Professional_2026.pptx"
    prs.save(pptx_path)
    return pptx_path

if __name__ == "__main__":
    path = create_professional_presentation()
    print(f"Professional PPTX generated at: {path}")
